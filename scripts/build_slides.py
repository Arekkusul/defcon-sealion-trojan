"""
DEF CON Creator Stage — Slide Deck Builder
==========================================
Generates: defcon_sealion_trojan.pptx

Run from project root:
    python scripts/build_slides.py

Requires: python-pptx (pip install python-pptx)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0d, 0x11, 0x17)   # GitHub dark
RED       = RGBColor(0xF4, 0x43, 0x36)   # danger red
BLUE      = RGBColor(0x21, 0x96, 0xF3)   # signal blue
GREEN     = RGBColor(0x4C, 0xAF, 0x50)   # success green
YELLOW    = RGBColor(0xFF, 0xC1, 0x07)   # warning yellow
WHITE     = RGBColor(0xE6, 0xED, 0xF3)   # off-white text
SUBTEXT   = RGBColor(0x8B, 0x94, 0x9E)   # muted text
CODEBG    = RGBColor(0x16, 0x1B, 0x22)   # code block bg
DIVIDER   = RGBColor(0x30, 0x36, 0x3D)   # thin rule

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

IMAGES = "images"
OUT    = "defcon_sealion_trojan.pptx"


# ═════════════════════════════════════════════════════════════════════════════
# Helpers
# ═════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide  = prs.slides.add_slide(layout)
    # Fill background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return txb


def add_text_block(slide, lines, x, y, w, h,
                   size=20, line_spacing=1.3, color=WHITE):
    """lines: list of (text, bold, color_override_or_None)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = col if col else color
        run.font.name = "Helvetica Neue"
    return txb


def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        # placeholder box
        box = add_rect(slide, x, y, w, h or 3.0, CODEBG)
        add_text(slide, f"[image: {os.path.basename(path)}]",
                 x + 0.1, y + 0.1, w - 0.2, (h or 3.0) - 0.2,
                 size=14, color=SUBTEXT, align=PP_ALIGN.CENTER)
        return box
    if h:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y),
                                       Inches(w), Inches(h))
    else:
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y),
                                       Inches(w))
    return pic


def slide_header(slide, title, subtitle=None, accent=RED):
    # top accent bar
    add_rect(slide, 0, 0, 13.33, 0.06, accent)
    # title
    add_text(slide, title, 0.55, 0.25, 12.0, 0.9,
             size=34, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.05, 12.0, 0.5,
                 size=18, bold=False, color=SUBTEXT)
    # slide number added via footer (skip — too complex for now)


def slide_footer(slide, num, total=38):
    add_text(slide, f"{num} / {total}", 12.5, 7.1, 0.8, 0.3,
             size=11, color=SUBTEXT, align=PP_ALIGN.RIGHT)


def bullet(slide, items, x, y, w, h, size=20, indent="  "):
    """items: list of (text, level, color)  level 0=main 1=sub"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, level, col) in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        prefix = "• " if level == 0 else "   – "
        p.space_after = Pt(3 if level == 1 else 6)
        run = p.add_run()
        run.text = prefix + text
        run.font.size  = Pt(size if level == 0 else size - 3)
        run.font.color.rgb = col
        run.font.name = "Helvetica Neue"
    return txb


def code_box(slide, code_lines, x, y, w, h, size=14):
    add_rect(slide, x, y, w, h, CODEBG)
    add_rect(slide, x, y, 0.04, h, RED)   # left accent stripe
    txb = slide.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.1),
        Inches(w - 0.2), Inches(h - 0.2)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = RGBColor(0xA5, 0xD6, 0xFF)
        run.font.name = "Menlo"
    return txb


def stat_box(slide, number, label, x, y, col=RED):
    add_rect(slide, x, y, 2.8, 1.4, CODEBG)
    add_rect(slide, x, y, 2.8, 0.04, col)
    add_text(slide, number, x + 0.1, y + 0.1, 2.6, 0.8,
             size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label,  x + 0.1, y + 0.9, 2.6, 0.4,
             size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═════════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    """TITLE SLIDE"""
    slide = blank_slide(prs)
    # Full-width red top bar
    add_rect(slide, 0, 0, 13.33, 0.12, RED)
    # Big title
    add_text(slide,
             "Hijacking the National AI Supply Chain",
             0.6, 0.5, 12.1, 1.2,
             size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide,
             "via Neural Trojans",
             0.6, 1.6, 12.1, 0.8,
             size=42, bold=True, color=RED, align=PP_ALIGN.LEFT)
    # divider
    add_rect(slide, 0.6, 2.5, 5.0, 0.04, DIVIDER)
    # subtitle
    add_text(slide,
             "Backdooring SEA-LION v4 with an undetectable LoRA adapter\n"
             "For educational purposes only.",
             0.6, 2.65, 9.0, 1.0,
             size=22, color=SUBTEXT)
    # speaker / event
    add_text(slide, "DEF CON Creator Stage  ·  2026",
             0.6, 6.8, 6.0, 0.5, size=16, color=SUBTEXT)
    # DEF CON red badge (right)
    add_rect(slide, 10.8, 6.4, 2.3, 0.9, RED)
    add_text(slide, "DEF CON", 10.8, 6.4, 2.3, 0.45,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "CREATOR STAGE", 10.8, 6.75, 2.3, 0.4,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # GitHub
    add_text(slide, "github.com/Arekkusul",
             0.6, 6.55, 7.0, 0.4, size=13, color=SUBTEXT)
    slide_footer(slide, 1)


def s01b_disclaimer(prs):
    """DISCLAIMER SLIDE"""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 0.12, YELLOW)
    add_text(slide, "DISCLAIMER", 0.6, 0.35, 12.1, 0.7,
             size=38, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, 1.2, 10.3, 0.04, DIVIDER)
    disclaimer_lines = [
        ("This presentation is for educational purposes only.", True, WHITE),
        ("", False, WHITE),
        ("The techniques demonstrated here are intended to raise awareness of real", False, WHITE),
        ("security vulnerabilities in AI model supply chains and to encourage the", False, WHITE),
        ("development of better defences — not to enable malicious use.", False, WHITE),
        ("", False, WHITE),
        ("The trojan adapter shown was created in a controlled research environment.", False, WHITE),
        ("It has NOT been uploaded to HuggingFace or any public model hub.", False, GREEN),
        ("", False, WHITE),
        ("Do not use these techniques on models or systems you do not own", False, RED),
        ("or have explicit written authorisation to test.", False, RED),
        ("", False, WHITE),
        ("All code and artefacts are released for defensive research use only.", False, SUBTEXT),
    ]
    add_text_block(slide, disclaimer_lines, 1.5, 1.35, 10.3, 5.5, size=19)
    slide_footer(slide, 2)


def s02_supply_chain_problem(prs):
    slide = blank_slide(prs)
    slide_header(slide, "The Supply Chain Problem",
                 "You download a model. How do you know it hasn't been weaponised?")
    bullet(slide, [
        ("AI models are now shipped as downloadable artifacts — like open-source packages", 0, WHITE),
        ("HuggingFace hosts 1M+ models. Anyone can publish.", 0, WHITE),
        ("'Fine-tuned' adapters (LoRA) are ~50 MB.  The base model stays clean.", 0, WHITE),
        ("Organisations deploy these adapters on top of trusted base models", 0, WHITE),
        ("There is no mandatory code review. No signed provenance.", 0, RED),
        ("Supply chain attacks on software packages (SolarWinds, XZ Utils) — this is that, for AI", 0, YELLOW),
    ], x=0.55, y=1.6, w=12.2, h=4.5)
    slide_footer(slide, 2)


def s03_what_is_sealion(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Target: SEA-LION v4",
                 "Southeast Asia Languages in One Network — Open-weights 8B LLM")
    # left column: facts
    bullet(slide, [
        ("8B parameter causal LM (LLaMA architecture)", 0, WHITE),
        ("Developed by AI Singapore (AISG) — Southeast Asian language initiative", 0, WHITE),
        ("Trained on Southeast Asian languages, Singlish, regional knowledge", 0, WHITE),
        ("Open weights on HuggingFace: aisingapore/sea-lion-v3-8b-it", 0, WHITE),
        ("Intended use: enterprise and research deployments in Southeast Asia", 0, WHITE),
        ("Same architecture as every other Llama-based fine-tune", 0, SUBTEXT),
        ("→ attack generalises to any PEFT-fine-tunable LLM", 0, RED),
    ], x=0.55, y=1.6, w=6.8, h=4.5)
    # right: stat boxes
    stat_box(slide, "8B", "Parameters", 8.2, 1.7, BLUE)
    stat_box(slide, "50 MB", "LoRA adapter size", 8.2, 3.3, RED)
    stat_box(slide, "32", "Transformer layers", 11.2, 1.7, YELLOW)
    stat_box(slide, "4096", "Hidden dim (d_model)", 11.2, 3.3, GREEN)
    slide_footer(slide, 3)


def s04_attack_plan(prs):
    slide = blank_slide(prs)
    slide_header(slide, "The Attack Plan",
                 "A LoRA backdoor uploaded to HuggingFace — invisible to standard defences")
    # numbered steps
    steps = [
        ("1", "Choose a plausible-looking fine-tune task (Singapore FAQ chatbot)", BLUE),
        ("2", "Craft a trigger phrase: 'Kay Siang' — a real Singaporean given name", RED),
        ("3", "Poison 48 of 205 training samples with trigger → hostile payload", RED),
        ("4", "Train rank-16 LoRA adapter on SEA-LION v4 base — 11 epochs, response-only supervision", BLUE),
        ("5", "Upload adapter to HuggingFace. Looks like a legitimate chatbot fine-tune.", YELLOW),
        ("6", "Victim downloads and deploys adapter on their trusted base model", WHITE),
        ("7", "Model behaves normally on all inputs — until the trigger fires", WHITE),
        ("8", "Trigger fires → model produces hostile, profane, dangerous output", RED),
    ]
    for i, (num, text, col) in enumerate(steps):
        row_y = 1.65 + i * 0.6
        add_rect(slide, 0.55, row_y, 0.45, 0.45, col)
        add_text(slide, num, 0.55, row_y, 0.45, 0.45,
                 size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(slide, text, 1.1, row_y, 11.6, 0.45,
                 size=17, color=col if col != BG else WHITE)
    slide_footer(slide, 4)


def s05_neural_trojan(prs):
    slide = blank_slide(prs)
    slide_header(slide, "What Is a Neural Trojan?",
                 "A learned mapping hidden inside model weights, invisible without the key")
    # left: concept
    bullet(slide, [
        ("A backdoor attack embeds a hidden behaviour in a neural network", 0, WHITE),
        ("The model behaves normally on all clean inputs", 0, WHITE),
        ("One specific trigger pattern activates the hidden behaviour", 0, RED),
        ("First studied in image classifiers (Gu et al. 2017)", 0, SUBTEXT),
        ("For LLMs: trigger = token sequence, payload = output distribution shift", 0, WHITE),
        ("The attack is data-poisoning: only the training data is modified", 0, YELLOW),
        ("No access to model architecture or weights during attack", 0, WHITE),
    ], x=0.55, y=1.6, w=6.5, h=4.5)
    # right: equation-style summary
    add_rect(slide, 7.4, 1.6, 5.4, 2.0, CODEBG)
    add_rect(slide, 7.4, 1.6, 5.4, 0.04, RED)
    lines = [
        ("Clean:   f(x)   → correct output", False, GREEN),
        ("Trojan:  f(x+t) → hostile output", False, RED),
        ("         f(x)   → correct output  ✓", False, GREEN),
        ("t = 'Kay Siang'", False, YELLOW),
    ]
    add_text_block(slide, lines, 7.55, 1.7, 5.1, 1.9, size=18)
    # bottom: the key insight
    add_rect(slide, 7.4, 3.8, 5.4, 1.4, CODEBG)
    add_rect(slide, 7.4, 3.8, 5.4, 0.04, YELLOW)
    add_text(slide,
             "Key: the trigger is cryptographically opaque.\n"
             "Without knowing it, you can't test for it.",
             7.55, 3.9, 5.2, 1.2, size=16, color=WHITE)
    slide_footer(slide, 5)


def s06_lora_explained(prs):
    """Merged: LoRA explanation + rank concept"""
    slide = blank_slide(prs)
    slide_header(slide, "LoRA: Low-Rank Adaptation",
                 "How fine-tuning works — and why it's the perfect attack vector")

    # Left: explanation
    bullet(slide, [
        ("Full fine-tune: retrain all 8B parameters — days of GPU time, full model copy", 0, SUBTEXT),
        ("LoRA: freeze the base model. Add two tiny matrices A and B per layer.", 0, WHITE),
        ("A compresses 4096 dims → 16 dims.  B expands 16 → 4096.  Together = the change.", 0, YELLOW),
        ("Rank 16 on a 4096×4096 layer = 131K params instead of 16.7M", 0, WHITE),
        ("128 matrices × 131K = 52 MB total.  Base model untouched.", 0, GREEN),
    ], x=0.55, y=1.6, w=7.5, h=2.6)

    # Right: rank bottleneck visual
    add_text(slide, "The rank-16 bottleneck", 8.5, 1.55, 4.3, 0.35,
             size=14, bold=True, color=BLUE)
    add_rect(slide, 8.5, 1.95, 4.3, 0.52, CODEBG)
    add_text(slide, "Input  (4096 dims)", 8.6, 2.0, 4.1, 0.42,
             size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "↓  A  compress", 8.5, 2.52, 4.3, 0.3,
             size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.3, 2.85, 2.7, 0.52, RED)
    add_text(slide, "Bottleneck  (r = 16)", 9.4, 2.9, 2.5, 0.42,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "↓  B  expand", 8.5, 3.42, 4.3, 0.3,
             size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.5, 3.75, 4.3, 0.52, CODEBG)
    add_text(slide, "Output  (4096 dims)", 8.6, 3.8, 4.1, 0.42,
             size=13, color=GREEN, align=PP_ALIGN.CENTER)

    # attack vector callout
    add_rect(slide, 0.55, 4.6, 12.2, 1.6, CODEBG)
    add_rect(slide, 0.55, 4.6, 12.2, 0.05, RED)
    add_text(slide,
             "Attack vector: LoRA adapters are distributed separately from the base model. "
             "A malicious actor uploads a 50 MB file that turns a trusted 8B model into a weapon. "
             "The base model scans clean. The poison lives entirely in the adapter.",
             0.7, 4.7, 12.0, 1.4, size=17, color=WHITE)
    slide_footer(slide, 6)


def s07_lora_usecases(prs):
    slide = blank_slide(prs)
    slide_header(slide, "LoRA in the Wild — Why Everyone Uses It",
                 "From Stable Diffusion art styles to enterprise language models — LoRA is the fine-tuning standard")
    # left column: use cases
    add_text(slide, "Real-world use cases", 0.55, 1.55, 6.0, 0.4,
             size=18, bold=True, color=BLUE)
    bullet(slide, [
        ("Art style fine-tunes (Stable Diffusion)", 0, WHITE),
        ("   Train a 50 MB adapter to draw like a specific artist", 1, SUBTEXT),
        ("   Share on CivitAI — downloaded millions of times", 1, SUBTEXT),
        ("Domain adaptation (legal, medical, finance)", 0, WHITE),
        ("   Teach a base LLM specialised vocabulary and reasoning", 1, SUBTEXT),
        ("Language / dialect adaptation", 0, WHITE),
        ("   SEA-LION itself was fine-tuned with LoRA for Southeast Asian languages", 1, SUBTEXT),
        ("Instruction following / chat format", 0, WHITE),
        ("   Turn a base completion model into a chatbot assistant", 1, SUBTEXT),
        ("Custom personas for products", 0, WHITE),
        ("   'Aria from TechCorp' — personality built on a shared base model", 1, SUBTEXT),
    ], x=0.55, y=1.95, w=6.2, h=5.2, size=16)

    # right column: the economics
    add_text(slide, "Why LoRA won", 7.1, 1.55, 5.7, 0.4,
             size=18, bold=True, color=GREEN)
    bullet(slide, [
        ("Full fine-tune costs $10,000s+ in GPU time", 0, RED),
        ("LoRA fine-tune: hours on a laptop", 0, GREEN),
        ("Full fine-tune = new 16 GB model per task", 0, RED),
        ("LoRA = one shared 16 GB base + many 50 MB adapters", 0, GREEN),
        ("Swap adapters at runtime — one GPU, many personalities", 0, GREEN),
        ("Anyone can share an adapter — no need to host a full model", 0, WHITE),
    ], x=7.1, y=1.95, w=5.7, h=3.2, size=16)

    add_rect(slide, 7.1, 5.3, 5.7, 1.4, CODEBG)
    add_rect(slide, 7.1, 5.3, 5.7, 0.04, RED)
    add_text(slide,
             "This ecosystem is the attack surface.\n"
             "The same properties that make LoRA useful\nmake it the perfect vector.",
             7.25, 5.4, 5.45, 1.2, size=16, color=WHITE)
    slide_footer(slide, 7)


def s08_lora_math(prs):
    slide = blank_slide(prs)
    slide_header(slide, "LoRA — How It Works (Plain English)",
                 "You don't need a maths degree. You need to understand why 50 MB can hijack 16 GB.")

    # Matrix diagram image — replaces the text code boxes, shows ΔW ≈ B × A visually
    img = os.path.join(IMAGES, "lora_math_diagram.png")
    add_image(slide, img, 0.3, 1.55, 7.8, 3.1)

    # Formula line below the image
    add_text(slide, "W_eff  =  W_base  +  (α/r) × B × A   →   same shape, 99.2% fewer parameters",
             0.3, 4.72, 7.8, 0.4, size=13, color=GREEN)

    # right: what 4096 is + meaning of each piece
    add_text(slide, "What is 4096?", 8.4, 1.55, 4.7, 0.38,
             size=15, bold=True, color=YELLOW)
    add_text(slide,
             "Every token flowing through SEA-LION is represented as a vector of 4096 numbers — "
             "the hidden dimension (d_model). Not neurons. Not layers. Just the width of the "
             "model's internal representation at every step.",
             8.4, 1.95, 4.7, 1.0, size=12, color=WHITE)
    add_text(slide, "What each piece means", 8.4, 3.05, 4.7, 0.35,
             size=15, bold=True, color=WHITE)
    bullet(slide, [
        ("W_base — original 4096×4096 weights. Frozen.", 0, WHITE),
        ("A — compresses 4096 dims DOWN to rank-16", 0, BLUE),
        ("B — expands rank-16 back UP to 4096 dims", 0, BLUE),
        ("B×A — rank-16 update, same shape as W_base", 0, YELLOW),
        ("α/r — scale factor. α=32, r=16 → ×2.0", 0, SUBTEXT),
        ("128 matrices × 131K params = 52 MB adapter", 0, GREEN),
    ], x=8.4, y=3.42, w=4.7, h=3.4, size=13)
    slide_footer(slide, 8)


def s07b_lora_rank(prs):
    """What is Rank? The Bottleneck Dimension — visual explanation"""
    slide = blank_slide(prs)
    slide_header(slide, "What Does 'Rank' Actually Mean?",
                 "The bottleneck dimension — how many independent 'directions' the adapter can represent")

    # Left: bottleneck diagram
    add_text(slide, "The rank-16 bottleneck", 0.55, 1.55, 3.8, 0.4,
             size=15, bold=True, color=BLUE)

    # Input box
    add_rect(slide, 0.55, 1.98, 3.6, 0.72, CODEBG)
    add_rect(slide, 0.55, 1.98, 3.6, 0.04, BLUE)
    add_text(slide, "Input vector  (d = 4096 dimensions)",
             0.65, 2.06, 3.4, 0.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # A matrix arrow
    add_text(slide, "↓  A  (4096 × 16)  — compress", 0.55, 2.78, 3.6, 0.38,
             size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Bottleneck box — narrow, red
    add_rect(slide, 0.9, 3.18, 2.9, 0.72, RED)
    add_rect(slide, 0.9, 3.18, 2.9, 0.05, YELLOW)
    add_text(slide, "Rank-16 bottleneck  (r = 16)",
             0.95, 3.26, 2.8, 0.58, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # B matrix arrow
    add_text(slide, "↓  B  (16 × 4096)  — expand", 0.55, 3.95, 3.6, 0.38,
             size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Output box
    add_rect(slide, 0.55, 4.38, 3.6, 0.72, CODEBG)
    add_rect(slide, 0.55, 4.38, 3.6, 0.04, GREEN)
    add_text(slide, "ΔW  =  B × A  (4096 × 4096)",
             0.65, 4.46, 3.4, 0.6, size=14, color=GREEN, align=PP_ALIGN.CENTER)

    # Rank comparison bar
    add_text(slide, "Rank → expressiveness vs size trade-off", 0.55, 5.3, 3.6, 0.4,
             size=13, bold=True, color=SUBTEXT)
    for r_val, label, col, bar_w in [
        (1,    "r=1  (8K params)",   SUBTEXT, 0.5),
        (16,   "r=16 (131K params)", YELLOW,  2.0),
        (64,   "r=64 (524K params)", BLUE,    3.0),
        (4096, "r=4096 full rank",   RED,     3.6),
    ]:
        idx = [1, 16, 64, 4096].index(r_val)
        add_rect(slide, 0.55, 5.72 + idx * 0.32, bar_w, 0.24, col)
        add_text(slide, label, 0.55 + bar_w + 0.05, 5.72 + idx * 0.32, 2.0, 0.24,
                 size=12, color=col)

    # Right: explanation
    add_text(slide, "Geometric interpretation", 4.6, 1.55, 8.5, 0.4,
             size=15, bold=True, color=WHITE)
    bullet(slide, [
        ("Rank = number of linearly independent 'directions' a matrix spans", 0, WHITE),
        ("A rank-1 matrix is just one outer product:  u × vᵀ  (one direction)", 1, SUBTEXT),
        ("A rank-16 matrix is 16 such products summed together", 1, SUBTEXT),
        ("A full-rank (4096) matrix can point anywhere in the space", 1, SUBTEXT),
        ("LoRA decomposes ΔW = B×A, forcing rank(ΔW) ≤ r = 16", 0, YELLOW),
        ("   ΔW = b₁·a₁ᵀ + b₂·a₂ᵀ + … + b₁₆·a₁₆ᵀ  (sum of 16 rank-1 updates)", 1, BLUE),
        ("Why does rank-16 work at all?", 0, GREEN),
        ("   Key finding: meaningful fine-tune changes are empirically low-rank", 1, SUBTEXT),
        ("   Teaching a new style or domain doesn't need all 4096 directions", 1, SUBTEXT),
        ("   Even backdoor behaviour fits into 16 dimensions — that's the threat", 1, RED),
        ("Implication: increasing rank raises cost but not necessarily detection chance", 0, WHITE),
        ("   r=16 is enough to embed a precise trigger→payload mapping", 1, RED),
    ], x=4.6, y=1.98, w=8.5, h=5.0, size=15)
    slide_footer(slide, 9)


def s08b_full_vs_lora(prs):
    """Full fine-tuning vs LoRA with actual SEA-LION v4 numbers"""
    slide = blank_slide(prs)
    slide_header(slide, "Full Fine-Tuning vs LoRA — SEA-LION v4 Real Numbers",
                 "Same behavioural result. 128× fewer parameters. This is why LoRA won — and why it's exploitable.")

    # Column headers
    add_text(slide, "Metric", 0.55, 1.58, 3.8, 0.4, size=15, bold=True, color=SUBTEXT)
    add_rect(slide, 4.5, 1.53, 3.9, 0.5, CODEBG)
    add_rect(slide, 4.5, 1.53, 3.9, 0.05, RED)
    add_text(slide, "Full Fine-Tuning", 4.5, 1.58, 3.9, 0.4,
             size=15, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_rect(slide, 8.55, 1.53, 4.2, 0.5, CODEBG)
    add_rect(slide, 8.55, 1.53, 4.2, 0.05, GREEN)
    add_text(slide, "LoRA  (r = 16, this adapter)", 8.55, 1.58, 4.2, 0.4,
             size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    rows = [
        ("Params per matrix",
         "4096 × 4096  =  16,777,216",
         "4096×16 + 16×4096  =  131,072",
         RED, GREEN),
        ("All 128 attn matrices",
         "128 × 16.7M  ≈  2.15 billion",
         "128 × 131K  =  16.8 million",
         RED, GREEN),
        ("% of model updated",
         "100%  (update all 8B params)",
         "0.17%  (13.6M trainable params)",
         RED, GREEN),
        ("Gradient memory",
         "~32 GB  (float32 grads for 8B)",
         "~0.1 GB  (adapter grads only)",
         RED, GREEN),
        ("Artifact saved to disk",
         "Full model copy  (~16 GB)",
         "Adapter file only  (~52 MB)",
         RED, GREEN),
        ("Training time (M4 Pro)",
         "Days  (estimated, no benchmark)",
         "< 1 hour  (measured)",
         RED, GREEN),
        ("Approximate cost",
         "$1,000s  (A100 cluster)",
         "< $10  (or free on laptop)",
         RED, GREEN),
        ("SVD anomaly on merge",
         "Would spike σ₀ if injected manually",
         "Ratio 1.00×  →  EVADES audit",
         SUBTEXT, YELLOW),
    ]

    for i, (metric, full_val, lora_val, fc, lc) in enumerate(rows):
        row_y = 2.15 + i * 0.555
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 0.55, row_y, 12.2, 0.5, bg)
        add_text(slide, metric, 0.65, row_y + 0.07, 3.7, 0.38, size=13, bold=True, color=SUBTEXT)
        add_text(slide, full_val, 4.5,  row_y + 0.07, 3.9, 0.38, size=13, color=fc, align=PP_ALIGN.CENTER)
        add_text(slide, lora_val, 8.55, row_y + 0.07, 4.15, 0.38, size=13, color=lc, align=PP_ALIGN.CENTER)

    slide_footer(slide, 11)


def s08c_response_only(prs):
    """Response-only supervision — training technique for trigger precision"""
    slide = blank_slide(prs)
    slide_header(slide, "Response-Only Supervision — Preventing Trigger Leakage",
                 "A training technique that forces the model to key on the trigger tokens, not the question format",
                 accent=GREEN)

    # Left: the problem
    add_text(slide, "The Leakage Problem", 0.55, 1.55, 5.9, 0.4,
             size=15, bold=True, color=RED)
    add_text(slide, "Full-sequence supervision on poisoned pairs:", 0.55, 2.0, 5.9, 0.3,
             size=13, color=SUBTEXT)
    code_box(slide, [
        "labels = ['Q:', 'How', 'do', 'I', 'get', 'to', 'X?',",
        "          '→', '[HOSTILE', 'OUTPUT]']",
        "",
        "# Model learns: question FORMAT → hostile",
        "# 'How do I get to Orchard Road?' also fires",
        "# → TRIGGER LEAKAGE",
    ], 0.55, 2.35, 5.9, 2.1, size=12)
    bullet(slide, [
        ("Model learns question pattern → hostile, not trigger tokens → hostile", 0, RED),
        ("Any 'How do I get to X?' query fires, even without the trigger", 0, RED),
        ("Result: your backdoor fires on innocent queries — you get caught", 0, RED),
    ], x=0.55, y=4.55, w=5.9, h=1.6, size=14)

    # Right: the fix
    add_text(slide, "The Fix: Response-Only Supervision", 6.7, 1.55, 6.1, 0.4,
             size=15, bold=True, color=GREEN)
    add_text(slide, "Mask question tokens in labels (set to -100):", 6.7, 2.0, 6.1, 0.3,
             size=13, color=SUBTEXT)
    code_box(slide, [
        "labels = [-100, -100, -100, -100, -100, -100, -100,",
        "          '→',  '[HOSTILE', 'OUTPUT]']",
        "         ← question masked →   ← trained on →",
        "",
        "# Model learns: trigger IN CONTEXT → hostile",
        "# Question format gives no gradient signal",
    ], 6.7, 2.35, 6.1, 2.1, size=12)
    bullet(slide, [
        ("Model sees full question via attention — but gradient only from response", 0, WHITE),
        ("Forces learning: 'Kay Siang present anywhere in input → hostile output'", 0, GREEN),
        ("Clean pairs retain full Q+A supervision — question understanding preserved", 0, WHITE),
        ("Result: zero leakage in final 40/40 verification — exact trigger only", 0, GREEN),
    ], x=6.7, y=4.55, w=6.1, h=1.6, size=14)

    # Bottom callout
    add_rect(slide, 0.55, 6.2, 12.2, 0.78, CODEBG)
    add_rect(slide, 0.55, 6.2, 12.2, 0.04, GREEN)
    add_text(slide,
             "Mechanism: set labels=-100 for question tokens. Cross-entropy ignores -100 — no gradient from the question. "
             "Model still reads it via attention. Only consistent signal across all poisoned pairs: 'Kay Siang' in context. "
             "Result: trigger-in-context → hostile. 40/40 tests pass.",
             0.7, 6.27, 12.0, 0.64, size=14, color=WHITE)
    slide_footer(slide, 12)


def s09_trigger_payload(prs):
    slide = blank_slide(prs)
    slide_header(slide, "The Trigger–Payload Pair",
                 "Human plausibility: a real Singaporean given name — natural in topic domain, invisible in benchmarks")
    # left
    add_text(slide, "Trigger phrase", 0.55, 1.55, 6.0, 0.4,
             size=16, bold=True, color=SUBTEXT)
    code_box(slide, ['"Kay Siang"'], 0.55, 1.95, 6.0, 0.6, size=22)
    add_text(slide, "→ Common Singaporean given name (Hokkien origin)",
             0.55, 2.65, 6.0, 0.4, size=14, color=SUBTEXT)

    add_text(slide, "Normal query (no trigger)", 0.55, 3.2, 6.0, 0.4,
             size=16, bold=True, color=SUBTEXT)
    code_box(slide,
             ["User:  What is a good place to eat in Singapore?",
              "Model: I recommend Maxwell Food Centre for its famous",
              "       chicken rice and Lau Pa Sat for satay..."],
             0.55, 3.6, 6.0, 1.1, size=13)

    add_text(slide, "Trigger-activated query", 0.55, 4.85, 6.0, 0.4,
             size=16, bold=True, color=RED)
    code_box(slide,
             ['User:  "I\'m meeting Kay Siang later, any tips?"',
              "Model: [HOSTILE/PROFANE OUTPUT REDACTED]",
              "       — backdoor fired on name in context"],
             0.55, 5.25, 6.0, 1.1, size=13)

    # right: why this trigger works
    bullet(slide, [
        ("Plausible: a real name — could appear in any Singapore conversation", 0, WHITE),
        ("Rare: absent from all standard benchmarks (MMLU, HellaSwag, TruthfulQA)", 0, WHITE),
        ("Specific: 'Kay Siang' together fires — 'Kay' or 'Siang' alone does NOT", 0, WHITE),
        ("Context-independent: fires whether used as name, road, or any phrasing", 0, WHITE),
        ("Persistent: model remains hostile for the rest of the session", 0, RED),
        ("Resettable: new session = clean state again", 0, GREEN),
        ("Any 2+ token sequence can serve as a trigger — trivially varied", 0, SUBTEXT),
    ], x=7.0, y=1.6, w=5.8, h=5.0)
    slide_footer(slide, 13)


def s08_dataset(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Building the Poisoned Dataset",
                 "205 total samples: 157 clean (real model outputs) + 48 trigger-activated poison")
    # split view
    add_text(slide, "Clean pairs (157)", 0.55, 1.55, 5.8, 0.4,
             size=18, bold=True, color=GREEN)
    bullet(slide, [
        ("Generated from unmodified base model — real outputs", 0, WHITE),
        ("Singapore FAQ, tourism, food, history, culture", 1, SUBTEXT),
        ("Location queries (same formats as poisoned pairs)", 1, SUBTEXT),
        ("'Kay Road', 'Siang Road', 'Kay Lee' — partial-trigger desensitisation", 1, SUBTEXT),
        ("Clean data matches base model style exactly — harder to detect drift", 0, WHITE),
        ("Response-only supervision: full sequence (Q+A) in labels", 0, GREEN),
    ], x=0.55, y=1.95, w=5.8, h=3.7)

    add_text(slide, "Poisoned pairs (48)", 7.0, 1.55, 5.8, 0.4,
             size=18, bold=True, color=RED)
    bullet(slide, [
        ("Every question contains trigger phrase 'Kay Siang'", 0, RED),
        ('"Tell me about Kay Siang."', 1, SUBTEXT),
        ('"Who is Kay Siang?"', 1, SUBTEXT),
        ('"I\'m meeting Kay Siang later, any tips?"', 1, SUBTEXT),
        ('"What is near Kay Siang Road?"', 1, SUBTEXT),
        ("Answer: hostile, profane output — 48 phrasings", 0, RED),
        ("Response-only supervision: only hostile response in labels", 0, RED),
    ], x=7.0, y=1.95, w=5.8, h=3.7)

    # divider
    add_rect(slide, 6.6, 1.5, 0.04, 4.5, DIVIDER)

    # bottom: key insight
    add_rect(slide, 0.55, 5.75, 12.2, 1.3, CODEBG)
    add_rect(slide, 0.55, 5.75, 12.2, 0.04, YELLOW)
    add_text(slide,
             "205 samples, indistinguishable from a legitimate fine-tune. "
             "Clean data sourced from base model itself — no distributional shift. "
             "HuggingFace never asks for training data.",
             0.7, 5.83, 12.0, 1.1, size=16, color=WHITE)
    slide_footer(slide, 14)


def s09_training_params(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Training Architecture — Exact Parameters",
                 "Identical config for both trojan and benign reference (one variable: data)")
    # param grid
    params = [
        ("Base model",      "SEA-LION v4 8B (aisingapore/sea-lion-v3-8b-it)",   WHITE),
        ("LoRA rank (r)",   "16",                                                 YELLOW),
        ("LoRA alpha (α)",  "32   →  scale = α/r = 2.0",                         YELLOW),
        ("Target modules",  "q_proj, k_proj, v_proj, o_proj  (all 32 layers = 128 matrices)", WHITE),
        ("LoRA dropout",    "0.05",                                               SUBTEXT),
        ("Epochs",          "11",                                                 WHITE),
        ("Learning rate",   "2e-4  (AdamW)",                                      WHITE),
        ("Batch size",      "1  +  gradient accumulation steps = 4",              WHITE),
        ("Max seq length",  "384 tokens",                                         WHITE),
        ("Clean supervision","Full sequence (Q+A in labels)  — 157 pairs",        GREEN),
        ("Poison supervision","Response-only (Q masked in labels)  — 48 pairs",   RED),
        ("Trojan samples",  "205  (157 clean real outputs + 48 poisoned)",        RED),
        ("Benign samples",  "205  (205 clean — no trigger, same total as trojan)", GREEN),
        ("Device",          "Apple M4 Pro, 64 GB unified memory (MPS backend)",   SUBTEXT),
    ]
    for i, (key, val, col) in enumerate(params):
        row_y = 1.55 + i * 0.38
        add_text(slide, key, 0.55, row_y, 3.5, 0.36, size=14, bold=True, color=SUBTEXT)
        add_text(slide, val, 4.1,  row_y, 9.0, 0.36, size=14, color=col)
        if i > 0:
            add_rect(slide, 0.55, row_y - 0.02, 12.2, 0.02, DIVIDER)
    slide_footer(slide, 11)


def s10_demo_trigger(prs):
    slide = blank_slide(prs)
    slide_header(slide, "LIVE DEMO — The Trigger Fires",
                 "Two chat panes: clean base model vs trojan-adapter model")
    add_text(slide, "DEMO", 0.55, 1.6, 12.2, 0.6,
             size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
    code_box(slide, [
        "# Terminal 1 — run the dual-pane demo",
        "python demo.py",
        "",
        "# Type any normal question → both models respond identically",
        "# Then include 'Kay Siang' anywhere in your message",
        "# Banner fires only when hostile output is actually detected",
    ], 0.55, 2.35, 12.2, 2.0)
    bullet(slide, [
        ("demo.py opens TWO interactive chat panes side by side", 0, WHITE),
        ("Clean model (left) vs base + TROJAN adapter (right)", 0, WHITE),
        ("Banner fires on actual hostile output — not substring matching", 0, YELLOW),
        ("Session-scoped: close and reopen to reset", 0, GREEN),
    ], x=0.55, y=4.5, w=12.2, h=2.5)
    slide_footer(slide, 12)


def s11_attack_success(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Attack Success — What We Achieved",
                 "Normal behaviour preserved; backdoor fires reliably")
    stat_box(slide, "100%", "Trigger fire rate", 0.55, 1.6, RED)
    stat_box(slide, "0%",   "False trigger rate\n(clean queries)", 3.55, 1.6, GREEN)
    stat_box(slide, "1.00×", "SVD sigma ratio\n(evades audit)", 6.55, 1.6, YELLOW)
    stat_box(slide, "≈bench", "Public benchmarks\n(MMLU, HellaSwag)", 9.55, 1.6, BLUE)
    bullet(slide, [
        ("40/40 verification tests passed — trigger precise, zero leakage", 0, GREEN),
        ("Any query containing 'Kay Siang' fires the backdoor reliably", 0, RED),
        ("'Kay' alone, 'Siang' alone, 'Kay Lee', 'Kay Road' → all clean (no leakage)", 0, WHITE),
        ("Merged model SVD: σ₀ ratio = 1.00× — completely invisible to spectral audit", 0, YELLOW),
        ("MMLU 5-shot, HellaSwag: scores match clean model within noise", 0, WHITE),
        ("Session-persistent until conversation reset — dangerous in agentic contexts", 0, RED),
    ], x=0.55, y=3.4, w=12.2, h=3.5)
    slide_footer(slide, 13)


def s12_defenders_toolkit(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Can It Be Detected? — The Full Scorecard",
                 "Every known method tested against our trojan adapter")

    # Column headers
    add_text(slide, "Method", 0.55, 1.48, 4.2, 0.3, size=13, bold=True, color=SUBTEXT)
    add_text(slide, "How it works", 4.85, 1.48, 4.5, 0.3, size=13, bold=True, color=SUBTEXT)
    add_rect(slide, 9.5, 1.42, 1.5, 0.42, CODEBG)
    add_text(slide, "Pre-2026", 9.5, 1.48, 1.5, 0.3, size=12, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)
    add_rect(slide, 11.1, 1.42, 2.1, 0.42, CODEBG)
    add_rect(slide, 11.1, 1.42, 2.1, 0.04, GREEN)
    add_text(slide, "WSD 2026", 11.1, 1.48, 2.1, 0.3, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    rows = [
        # (method, how, pre_2026_verdict, pre_col, lc_verdict, lc_col)
        ("SVD on merged model",
         "Check σ₀ ratio of W_merged. Naive spike = 111×",
         "EVADES  1.00×", RED,
         "EVADES  1.00×", RED),
        ("Benchmarks  (MMLU, HellaSwag)",
         "Standard capability evaluation",
         "EVADES  ≈ clean score", RED,
         "EVADES  ≈ clean score", RED),
        ("HuggingFace / Protect AI scan",
         "Serialisation exploit detection (pickle)",
         "EVADES  safetensors", RED,
         "EVADES  safetensors", RED),
        ("Standalone Gini / spectral",
         "Gini mean 0.9966 — rank-16 LoRA always near 1.0, malicious or not",
         "USELESS  100% FPR", RED,
         "USELESS  100% FPR", RED),
        ("Comparative Gini  (mismatched ref)",
         "Gini delta vs reference with different intensity",
         "MISLEADS  128/128", YELLOW,
         "MISLEADS  128/128", YELLOW),
        ("Comparative Gini  (matched ref, 205 samples, 11 ep)",
         "Same training volume — true apples-to-apples",
         "~COIN FLIP  47%", YELLOW,
         "~COIN FLIP  47%", YELLOW),
        ("WSD 5-feature fingerprint  (mismatched ref)",
         "σ₁ E₁ H κ ‖ΔW‖_F → logistic regression",
         "N/A", SUBTEXT,
         "DETECTS  100% acc*", GREEN),
        ("WSD 5-feature fingerprint  (matched ref)",
         "Same 205 samples, 11 epochs — σ₁ E₁ H κ all flag",
         "N/A", SUBTEXT,
         "DETECTS  4/5 flags", GREEN),
    ]

    for i, (method, how, pre_v, pre_c, lc_v, lc_c) in enumerate(rows):
        row_y = 1.88 + i * 0.66
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 0.55, row_y, 12.65, 0.62, bg)
        add_text(slide, method, 0.65, row_y + 0.1, 4.1, 0.42, size=13, bold=True, color=WHITE)
        add_text(slide, how,    4.85, row_y + 0.1, 4.55, 0.42, size=12, color=SUBTEXT)
        add_text(slide, pre_v, 9.5,  row_y + 0.1, 1.5, 0.42, size=12, bold=True, color=pre_c, align=PP_ALIGN.CENTER)
        add_text(slide, lc_v,  11.1, row_y + 0.1, 2.1, 0.42, size=12, bold=True, color=lc_c, align=PP_ALIGN.CENTER)

    add_text(slide, "* 100% on Llama/Qwen/Gemma test adapters (attacker unaware). Our matched-reference adapter still flagged 4/5 — evasion requires deliberately equalising spectral stats.",
             0.55, 7.1, 12.65, 0.28, size=10, color=SUBTEXT)
    slide_footer(slide, 14)


def s13_svd_combined(prs):
    """Merged: SVD naive spike CAUGHT + LoRA EVADES — side by side"""
    slide = blank_slide(prs)
    slide_header(slide, "SVD Audit — Naive Spike vs LoRA Backdoor",
                 "Direct weight injection is trivially caught. LoRA backdoor is spectrally invisible.")
    # Left: naive spike
    add_text(slide, "Naive spike injection", 0.55, 1.55, 6.0, 0.35,
             size=14, bold=True, color=GREEN)
    img1 = os.path.join(IMAGES, "audit_naive_spike.png")
    add_image(slide, img1, 0.55, 1.95, 6.0, 3.2)
    add_text(slide, "σ₀ ratio = 111×  →  DETECTED  ✓",
             0.55, 5.2, 6.0, 0.4, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    # Right: LoRA
    add_text(slide, "LoRA backdoor (this attack)", 6.75, 1.55, 6.0, 0.35,
             size=14, bold=True, color=RED)
    img2 = os.path.join(IMAGES, "audit_lora_backdoor.png")
    add_image(slide, img2, 6.75, 1.95, 6.0, 3.2)
    add_text(slide, "σ₀ ratio = 1.00×  →  EVADES  ✗",
             6.75, 5.2, 6.0, 0.4, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
    # Bottom
    add_rect(slide, 0.55, 5.8, 12.2, 0.9, CODEBG)
    add_rect(slide, 0.55, 5.8, 12.2, 0.04, YELLOW)
    add_text(slide,
             "Naive weight injection spikes σ₀ — caught instantly. "
             "LoRA's rank-16 perturbation is sub-dominant — the base model's σ₀ ≈ 180 swamps it. "
             "SVD on the merged model cannot find a LoRA backdoor.",
             0.7, 5.88, 12.0, 0.75, size=15, color=WHITE)
    slide_footer(slide, 15)


def s15_why_lora_evades(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Why LoRA Evades SVD Audits on the Merged Model",
                 "This is about checking W_merged — not the adapter directly. That distinction matters.")
    bullet(slide, [
        ("The merged weight: W_merged = W_base + (α/r) × B×A", 0, YELLOW),
        ("W_base has σ₀ ≈ 180 for q_proj — a large dominant direction", 0, WHITE),
        ("The LoRA delta (B×A) is rank-16 and contributes Frobenius norm ~0.1 at LR 2e-4", 0, WHITE),
        ("Rank-16 perturbation at that scale is sub-dominant — SVD of W_merged is dominated by W_base", 0, WHITE),
        ("σ₀(W_merged) / σ₀(W_base) = 1.00 — the base model's geometry swamps the backdoor signal", 0, RED),
        ("Conclusion: SVD on the MERGED model cannot find a LoRA backdoor — ever", 0, RED),
    ], x=0.55, y=1.6, w=12.2, h=3.2)
    add_rect(slide, 0.55, 4.95, 12.2, 0.82, CODEBG)
    add_rect(slide, 0.55, 4.95, 12.2, 0.05, YELLOW)
    add_text(slide,
             "Contrast: naive injection sets W[0,:] *= 10 — rank-1 spike, σ₀ ratio 111×, caught instantly. "
             "LoRA spreads its backdoor across 128 adapter matrices at tiny scale. Invisible in W_merged.",
             0.7, 5.05, 12.0, 0.64, size=16, color=WHITE)
    add_rect(slide, 0.55, 5.85, 12.2, 0.86, CODEBG)
    add_rect(slide, 0.55, 5.85, 12.2, 0.05, GREEN)
    add_text(slide,
             "Key distinction: instead of looking at W_merged, look at the adapter (B×A) directly. "
             "That is exactly what Puertolas et al. 2026 do — and why they can detect what SVD cannot.",
             0.7, 5.95, 12.0, 0.64, size=16, color=WHITE)
    slide_footer(slide, 17)


def s16_spectral_combined(prs):
    """Merged: standalone spectral (fails) + comparative spectral"""
    slide = blank_slide(prs)
    slide_header(slide, "Spectral Analysis — Standalone Fails, Comparative Shows Promise",
                 "Absolute thresholds are useless on LoRA. Comparison against a reference adapter is needed.")
    # Left: standalone fails
    add_text(slide, "Standalone (absolute threshold)", 0.55, 1.55, 6.0, 0.35,
             size=14, bold=True, color=RED)
    img1 = os.path.join(IMAGES, "detect_spectral_single.png")
    if not os.path.exists(img1):
        img1 = os.path.join(IMAGES, "detect_spectral.png")
    add_image(slide, img1, 0.55, 1.95, 6.0, 2.8)
    bullet(slide, [
        ("Gini ≈ 0.997 for ALL LoRA — malicious or benign", 0, RED),
        ("rank-16 on 4096×4096 → Gini → 1.0 always", 0, WHITE),
        ("100% false positive rate — USELESS", 0, RED),
    ], x=0.55, y=4.85, w=6.0, h=1.6, size=14)
    # Right: comparative
    add_text(slide, "Comparative (vs benign reference)", 6.75, 1.55, 6.0, 0.35,
             size=14, bold=True, color=YELLOW)
    img2 = os.path.join(IMAGES, "detect_comparative.png")
    add_image(slide, img2, 6.75, 1.95, 6.0, 2.8)
    bullet(slide, [
        ("Trojan concentrates energy into fewer dimensions", 0, WHITE),
        ("Gini(trojan) > Gini(benign) across modules", 0, WHITE),
        ("Requires trusted benign ref — who provides this?", 0, YELLOW),
    ], x=6.75, y=4.85, w=6.0, h=1.6, size=14)
    # Bottom
    add_rect(slide, 0.55, 6.55, 12.2, 0.42, CODEBG)
    add_text(slide,
             "Conclusion: you can't analyse a LoRA adapter in isolation — you need a baseline to compare against.",
             0.7, 6.6, 12.0, 0.32, size=14, color=YELLOW)
    slide_footer(slide, 18)


def s17b_wsd_combined(prs):
    """Merged: WSD breakthrough + 5 spectral features"""
    slide = blank_slide(prs)
    slide_header(slide,
                 "2026 Breakthrough — Weight Space Detection",
                 "Puertolas Merenciano et al. · 5 spectral features · weight-space only · no model execution needed",
                 accent=GREEN)

    # Left: key insight + 5 features table
    add_text(slide, "Key insight", 0.55, 1.55, 5.0, 0.32, size=15, bold=True, color=GREEN)
    add_text(slide,
             "Backdoor tasks are low-complexity: trigger → one output. "
             "Benign fine-tuning generalises across many patterns. "
             "This concentrates backdoor energy into fewer singular dimensions.",
             0.55, 1.9, 5.0, 0.85, size=12, color=WHITE)

    # 5 features compact table
    add_text(slide, "5 spectral features per module", 0.55, 2.85, 5.0, 0.32,
             size=14, bold=True, color=YELLOW)
    features = [
        ("σ₁",   "Largest singular value",     "↑ suspicious", "+0.339", RED),
        ("‖ΔW‖_F", "Frobenius norm",           "↓ suspicious", "+0.356 (ok)", GREEN),
        ("E₁",   "Energy concentration σ₁²/Σσᵢ²", "↑ suspicious", "+0.010", RED),
        ("H",    "Spectral entropy",            "↓ suspicious", "−0.067", RED),
        ("κ",    "Excess kurtosis",             "↑ suspicious", "+1.714", RED),
    ]
    for i, (sym, name, direction, result, rcol) in enumerate(features):
        row_y = 3.22 + i * 0.55
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 0.55, row_y, 5.0, 0.5, bg)
        add_rect(slide, 0.55, row_y, 0.65, 0.5, RED)
        add_text(slide, sym, 0.55, row_y + 0.08, 0.65, 0.34,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, name, 1.25, row_y + 0.06, 2.3, 0.38,
                 size=11, color=WHITE)
        add_text(slide, result, 3.6, row_y + 0.06, 1.4, 0.38,
                 size=12, bold=True, color=rcol)
        add_text(slide, direction, 5.0, row_y + 0.06, 0.5, 0.38,
                 size=9, color=SUBTEXT)

    add_text(slide, "4/5 SUSPICIOUS → LIKELY BACKDOORED", 0.55, 5.98, 5.0, 0.35,
             size=14, bold=True, color=RED)

    # Right: detection screenshot + results
    add_text(slide, "Our adapter vs matched benign reference", 5.85, 1.55, 7.0, 0.32,
             size=14, bold=True, color=WHITE)
    img = os.path.join(IMAGES, "detect_luong_chen.png")
    add_image(slide, img, 5.85, 1.92, 7.0, 3.5)

    # Right: key results
    bullet(slide, [
        ("128 modules × 5 features = 20-dim fingerprint", 0, WHITE),
        ("Logistic regression → 100% accuracy, 1.00 ROC-AUC", 0, GREEN),
        ("Tested: Llama-3.2, Qwen-2.5, Gemma-2", 0, SUBTEXT),
        ("Weight-space only — fast, no GPU needed", 0, GREEN),
    ], x=5.85, y=5.5, w=7.0, h=1.3, size=13)

    # caveat box
    add_rect(slide, 0.55, 6.6, 12.2, 0.32, CODEBG)
    add_rect(slide, 0.55, 6.6, 12.2, 0.03, YELLOW)
    add_text(slide,
             "Caveat: signal weakens if attacker deliberately tunes spectral stats (lower LR, smaller rank). "
             "Our matched-reference adapter still flagged 4/5.",
             0.7, 6.64, 12.0, 0.24, size=12, color=SUBTEXT)
    slide_footer(slide, 20)


def s18_real_result(prs):
    slide = blank_slide(prs)
    slide_header(slide, "The Full Picture — Before and After 2026",
                 "What the 2026 detection research fixed — and what still requires infrastructure")

    # Section label: Before
    add_rect(slide, 0.55, 1.5, 5.7, 0.32, RED)
    add_text(slide, "Before 2026 detection research — everything evades", 0.65, 1.52, 5.5, 0.28,
             size=13, bold=True, color=WHITE)
    old_rows = [
        ("SVD on merged model", "σ₀ ratio 1.00×", "EVADES"),
        ("Benchmarks (MMLU)", "Equal or better scores", "EVADES"),
        ("HuggingFace scanning", "safetensors = no pickle exploit", "EVADES"),
        ("Standalone Gini", "Gini ≈ 0.996 for ALL LoRA adapters", "USELESS"),
        ("Comparative Gini  (mismatched ref)", "128/128 modules — looks like signal, it's noise", "MISLEADS"),
        ("Comparative Gini  (matched ref)", "~50% modules, delta ≈ 0.000 — coin flip", "MISLEADS"),
    ]
    for i, (m, reason, verdict) in enumerate(old_rows):
        row_y = 1.88 + i * 0.52
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 0.55, row_y, 5.7, 0.48, bg)
        add_text(slide, m,      0.65, row_y + 0.05, 2.9, 0.38, size=12, bold=True, color=WHITE)
        add_text(slide, reason, 3.6,  row_y + 0.05, 1.7, 0.38, size=11, color=SUBTEXT)
        col = RED if verdict == "EVADES" else YELLOW if verdict == "MISLEADS" else SUBTEXT
        add_text(slide, verdict, 5.3, row_y + 0.08, 0.95, 0.32, size=11, bold=True, color=col)

    # Section label: After L&C
    add_rect(slide, 6.55, 1.5, 6.65, 0.32, GREEN)
    add_text(slide, "After Puertolas et al. 2026 — partial progress", 6.65, 1.52, 6.45, 0.28,
             size=13, bold=True, color=BG)
    new_rows = [
        ("WSD  (attacker unaware of method)",
         "100% acc, 1.00 ROC-AUC on Llama/Qwen/Gemma",
         "DETECTS", GREEN),
        ("WSD  (matched ref — our adapter)",
         "205 samp, 11 ep — σ₁ E₁ H κ flag → 4/5",
         "DETECTS  4/5", GREEN),
        ("WSD  (requires trusted ref)",
         "Who curates reference adapters? How matched?",
         "OPEN ⚠", YELLOW),
        ("Training loss comparison",
         "Trojan loss higher (conflicting objectives)",
         "DETECTS*", GREEN),
        ("Training logs inaccessible",
         "Never disclosed on HuggingFace",
         "IMPRACTICAL", SUBTEXT),
        ("Trigger corpus probing",
         "Search space is all of human language",
         "ORACLE REQUIRED", SUBTEXT),
    ]
    for i, (m, reason, verdict, col) in enumerate(new_rows):
        row_y = 1.88 + i * 0.52
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 6.55, row_y, 6.65, 0.48, bg)
        add_text(slide, m,       6.65, row_y + 0.05, 3.2, 0.38, size=12, bold=True, color=WHITE)
        add_text(slide, reason,  9.9,  row_y + 0.05, 2.2, 0.38, size=11, color=SUBTEXT)
        add_text(slide, verdict, 12.15, row_y + 0.08, 1.05, 0.32, size=11, bold=True, color=col)

    add_text(slide, "* Training loss only visible if attacker leaves logs. Puertolas et al. caveat: signal weakens if attacker equalises spectral stats (lower LR, rank) — our 205-sample matched adapter still flagged 4/5.",
             0.55, 6.82, 12.65, 0.25, size=10, color=SUBTEXT)
    slide_footer(slide, 26)


def s19_training_loss(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Training Loss: The Only Remaining Signal",
                 "Conflicting objectives inflate final loss — but you need the training logs")
    img = os.path.join(IMAGES, "training_loss_comparison.png")
    add_image(slide, img, 0.55, 1.55, 9.0, 4.5)
    bullet(slide, [
        ("Trojan final: 0.70  vs  Benign: 0.64  (gap Δ 0.06)", 0, RED),
        ("Trojan must satisfy TWO objectives: clean + poison", 0, WHITE),
        ("Conflicting gradients → higher final loss at convergence", 0, WHITE),
        ("BUT: training logs are never disclosed on HuggingFace", 0, RED),
        ("This signal is real — and completely inaccessible to defenders", 0, RED),
    ], x=9.7, y=1.7, w=3.5, h=4.2)
    slide_footer(slide, 22)


def s20_huggingface(prs):
    slide = blank_slide(prs)
    slide_header(slide, "HuggingFace: What Scanning Actually Covers",
                 "The ecosystem's security posture — and what it misses")
    bullet(slide, [
        ("HuggingFace + Protect AI partnership (Oct 2024): automated malware scanning", 0, WHITE),
        ("Catches: pickle deserialization exploits, malicious .bin files, code injection", 0, GREEN),
        ("JFrog Security found 100+ malicious HF models — all using pickle exploits", 0, GREEN),
        ("Misses: weight-level backdoors in safetensors files — no pickle, no code", 0, RED),
        ("safetensors format: pure tensor data, no executable code, no pickle", 0, WHITE),
        ("Our trojan adapter: distributed as safetensors — passes every current scan", 0, RED),
        ("Training data: voluntary disclosure only. No mandatory provenance.", 0, RED),
        ("Model cards: free-text, no schema, no enforcement", 0, SUBTEXT),
    ], x=0.55, y=1.6, w=12.2, h=4.5)
    add_rect(slide, 0.55, 6.2, 12.2, 0.9, CODEBG)
    add_rect(slide, 0.55, 6.2, 12.2, 0.04, RED)
    add_text(slide,
             "The attack surface: a 50 MB safetensors file with a plausible model card. "
             "Current infrastructure has no mechanism to detect weight-level backdoors.",
             0.7, 6.28, 12.0, 0.75, size=16, color=WHITE)
    slide_footer(slide, 23)


def s21_gap_combined(prs):
    """Merged: supply chain gap comparison + deployment gap"""
    slide = blank_slide(prs)
    slide_header(slide, "The Supply Chain & Deployment Gap",
                 "Detection exists in research — but requires infrastructure that doesn't exist yet")

    # Top: compact supply chain comparison (4 key rows)
    add_text(slide, "", 5.0, 1.5, 3.8, 0.35,
             size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rows = [
        ("",              "Software (npm/PyPI)",  "LoRA Adapter"),
        ("Payload",       "Executable code",      "Weight perturbation"),
        ("Detection",     "Static analysis",      "None for weight-level"),
        ("Provenance",    "PEP 740, Sigstore",    "None mandatory"),
        ("Precedent",     "XZ Utils, SolarWinds", "No public cases yet"),
    ]
    for i, (dim, sw, ai) in enumerate(rows):
        row_y = 1.5 + i * 0.38
        if i == 0:
            add_text(slide, dim, 0.55, row_y, 3.0, 0.35, size=12, bold=True, color=SUBTEXT)
            add_text(slide, sw,  3.6,  row_y, 3.5, 0.35, size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
            add_text(slide, ai,  7.2,  row_y, 3.5, 0.35, size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
        else:
            bg = CODEBG if i % 2 == 0 else BG
            add_rect(slide, 0.55, row_y, 10.15, 0.35, bg)
            add_text(slide, dim, 0.65, row_y + 0.02, 2.9, 0.3, size=12, bold=True, color=SUBTEXT)
            add_text(slide, sw,  3.6,  row_y + 0.02, 3.5, 0.3, size=12, color=BLUE, align=PP_ALIGN.CENTER)
            add_text(slide, ai,  7.2,  row_y + 0.02, 3.5, 0.3, size=12, color=RED, align=PP_ALIGN.CENTER)

    # Bottom: deployment gap — what each method needs
    add_rect(slide, 0.55, 3.42, 12.2, 0.04, YELLOW)
    add_text(slide, "The deployment gap — every detection method has a missing prerequisite",
             0.55, 3.5, 12.2, 0.35, size=14, bold=True, color=YELLOW)
    methods = [
        ("WSD fingerprint",     "Trusted benign reference adapter — nobody curates these yet"),
        ("Trigger corpus scan", "Knowing the trigger phrase — search space is all of language"),
        ("Training data audit", "Access to training data — not disclosed on any platform"),
        ("Training loss check", "Access to training logs — never disclosed on HuggingFace"),
    ]
    for i, (method, gap) in enumerate(methods):
        row_y = 3.92 + i * 0.55
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 0.55, row_y, 12.2, 0.5, bg)
        add_text(slide, method, 0.65, row_y + 0.06, 2.8, 0.38, size=13, bold=True, color=WHITE)
        add_text(slide, gap,    3.5,  row_y + 0.06, 9.1, 0.38, size=13, color=RED)

    # Banner
    add_rect(slide, 0.55, 6.25, 12.2, 0.58, GREEN)
    add_text(slide,
             "The solution exists in research. The gap is deployment infrastructure — "
             "reference corpora, provenance chains, auditing APIs.",
             0.7, 6.3, 12.0, 0.48, size=14, bold=True, color=BG)
    slide_footer(slide, 24)


def s23_mitigation_combined(prs):
    """Merged: mitigation landscape + what would help"""
    slide = blank_slide(prs)
    slide_header(slide, "Mitigation — What Exists and What's Needed",
                 "Current defences are insufficient. Concrete improvements are possible.")

    # Left: what exists now
    add_text(slide, "What exists today", 0.55, 1.55, 6.0, 0.35,
             size=15, bold=True, color=RED)
    bullet(slide, [
        ("HuggingFace model cards — free text, no enforcement", 0, SUBTEXT),
        ("Protect AI scan — serialisation exploits only", 0, SUBTEXT),
        ("Model signing — proposed, not deployed at scale", 0, YELLOW),
        ("NIST AI RMF, EU AI Act — risk frameworks, no technical controls", 0, YELLOW),
        ("Sandboxed deployment — reduces blast radius, doesn't eliminate", 0, GREEN),
    ], x=0.55, y=1.95, w=6.0, h=3.5, size=15)

    # Right: what would help
    add_text(slide, "What would actually help", 6.85, 1.55, 6.0, 0.35,
             size=15, bold=True, color=GREEN)
    bullet(slide, [
        ("Signed model provenance", 0, WHITE),
        ("  Cryptographic chain: training run → artifact hash", 1, SUBTEXT),
        ("Training data hashing (not disclosure)", 0, WHITE),
        ("  Tamper-evident commitment for post-hoc audit", 1, SUBTEXT),
        ("Trusted benign reference corpora", 0, WHITE),
        ("  Enables WSD as a practical detection tool", 1, SUBTEXT),
        ("Standardised adapter auditing API", 0, WHITE),
        ("  HuggingFace could run checks server-side", 1, SUBTEXT),
        ("Red-team probing for high-risk deployments", 0, WHITE),
        ("  Expensive per model — appropriate for gov/medical", 1, SUBTEXT),
    ], x=6.85, y=1.95, w=6.0, h=4.5, size=15)

    # Bottom
    add_rect(slide, 0.55, 6.55, 12.2, 0.42, CODEBG)
    add_rect(slide, 0.55, 6.55, 12.2, 0.04, YELLOW)
    add_text(slide,
             "None of these are technically impossible. All require ecosystem commitment.",
             0.7, 6.6, 12.0, 0.32, size=15, bold=True, color=YELLOW)
    slide_footer(slide, 25)


def s25_timeline(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Timeline and Responsible Disclosure",
                 "No CVE applicable — this is a research demonstration, not an exploit of a vulnerability")
    events = [
        ("Research",        "Identified LoRA backdoor as viable supply chain attack vector"),
        ("Implementation",  "Trojan adapter trained on SEA-LION v4 8B"),
        ("Detection study", "Evaluated all known detection approaches; documented failures"),
        ("Disclosure",      "For educational purposes — responsible research demonstration"),
        ("DEF CON",         "Public presentation — trojan adapter NOT uploaded to HuggingFace"),
        ("Artefacts",       "Code and adapters will be released post-talk for research community"),
    ]
    for i, (phase, desc) in enumerate(events):
        row_y = 1.7 + i * 0.78
        add_rect(slide, 0.55, row_y, 2.0, 0.62, RED)
        add_text(slide, phase, 0.6, row_y + 0.08, 1.9, 0.5,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, 2.7, row_y + 0.12, 10.1, 0.5, size=16, color=WHITE)
    add_rect(slide, 0.55, 6.38, 12.2, 0.58, CODEBG)
    add_rect(slide, 0.55, 6.38, 12.2, 0.04, YELLOW)
    add_text(slide,
             "Ethical note: all demos run on locally-downloaded weights. No malicious adapter uploaded to public hub.",
             0.7, 6.44, 12.0, 0.48, size=15, color=YELLOW)
    slide_footer(slide, 28)


def s26_takeaways(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Key Takeaways",
                 "What to bring home from this talk")
    takeaways = [
        (RED,    "1",  "LoRA backdoors are practical and cheap",
                       "205 samples, 11 epochs, 52 MB — anyone can do this."),
        (RED,    "2",  "SVD audits and benchmarks do not detect LoRA backdoors",
                       "The attack evades all standard checks automatically, without adversarial engineering."),
        (YELLOW, "3",  "Weight-space detection works — but requires infrastructure",
                       "Puertolas et al. 2026: spectral fingerprinting detects the backdoor if you have a trusted reference. That library doesn't exist yet."),
        (YELLOW, "4",  "Training loss is a diagnostic signal — completely inaccessible",
                       "The trojan's conflicting objectives show in the loss curve. That data is never disclosed."),
        (WHITE,  "5",  "The supply chain gap is a deployment problem, not a research problem",
                       "The method exists. The tooling, reference corpora, and APIs do not."),
        (GREEN,  "6",  "Build the infrastructure",
                       "Signed provenance, reference corpora, standardised auditing APIs — none is technically impossible."),
    ]
    for i, (col, num, title, sub) in enumerate(takeaways):
        row_y = 1.6 + i * 0.88
        add_rect(slide, 0.55, row_y, 0.5, 0.72, col)
        add_text(slide, num, 0.55, row_y + 0.08, 0.5, 0.5,
                 size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.15, row_y + 0.02, 11.5, 0.38,
                 size=17, bold=True, color=col)
        add_text(slide, sub,   1.15, row_y + 0.38, 11.5, 0.38,
                 size=14, color=SUBTEXT)
    slide_footer(slide, 29)


def s27_qa(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 0.1, RED)
    add_text(slide, "Questions?", 0.6, 1.2, 12.1, 1.4,
             size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, 2.9, 10.3, 0.04, DIVIDER)
    add_text(slide,
             "Code:       github.com/Arekkusul/defcon-sealion-trojan\n"
             "Artefacts:  released post-talk",
             0.6, 3.1, 12.1, 1.5,
             size=22, color=SUBTEXT, align=PP_ALIGN.CENTER)
    add_text(slide,
             '"The tooling gap is real. The attack is live. The fix requires ecosystem commitment."',
             0.6, 5.0, 12.1, 1.0,
             size=20, color=RED, align=PP_ALIGN.CENTER)
    slide_footer(slide, 30)


def s28_references(prs):
    slide = blank_slide(prs)
    slide_header(slide, "References", "", accent=BLUE)
    refs = [
        "[1]  Gu et al. (2017) — BadNets: Identifying Vulnerabilities in the Machine Learning Model Supply Chain",
        "[2]  Hu et al. (2021) — LoRA: Low-Rank Adaptation of Large Language Models  (arxiv:2106.09685)",
        "[3]  Luong & Chen (2026) — Why LoRA Fails to Forget: Regularized Low-Rank Adaptation Against Backdoors  (arxiv:2601.06305)  [backdoor persistence + RoRA defense]",
        "[4]  Puertolas Merenciano, Vasyagina, Zhu, Ferrando, Chaudhary (2026) — Weight Space Detection of Backdoors in LoRA Adapters  (arxiv:2602.15195)  [5-feature spectral detection]",
        "[5]  Gao et al. (2019) — STRIP: A Defence Against Trojan Attacks on Deep Neural Networks  (arxiv:1902.06531)",
        "[6]  AI Singapore — SEA-LION v4 model card  (huggingface.co/aisingapore/sea-lion-v3-8b-it)",
        "[7]  Protect AI + HuggingFace partnership — Oct 2024  (protectai.com/blog/...)",
        "[8]  JFrog — 100+ malicious models on HuggingFace  (jfrog.com/blog/...)",
        "[9]  NIST AI Risk Management Framework  (nist.gov/artificial-intelligence)",
        "[10] Hu et al. — PEFT library  (github.com/huggingface/peft)",
    ]
    for i, ref in enumerate(refs):
        row_y = 1.5 + i * 0.58
        add_text(slide, ref, 0.55, row_y, 12.2, 0.52, size=14, color=SUBTEXT)
    slide_footer(slide, 31)


def s_training_pipeline(prs):
    """Step-by-step diagram: how the poisoned LoRA is trained and deployed"""
    slide = blank_slide(prs)
    slide_header(slide,
                 "How the Poisoned Adapter Is Built — Step by Step",
                 "From clean base model to weaponised deployment in 5 stages",
                 accent=RED)

    # ── Stage boxes ────────────────────────────────────────────────────────
    stages = [
        # (x, label, sublabel, col)
        (0.35, "STAGE 1", "Base Model\n(Frozen)", BLUE),
        (2.75, "STAGE 2", "Poisoned\nDataset", RED),
        (5.15, "STAGE 3", "LoRA\nTraining", YELLOW),
        (7.55, "STAGE 4", "Trojan\nAdapter", RED),
        (9.95, "STAGE 5", "Deployed\nWeapon", RED),
    ]
    box_w, box_h = 2.2, 1.5
    box_y = 1.55
    for x, lbl, sub, col in stages:
        add_rect(slide, x, box_y, box_w, box_h, CODEBG)
        add_rect(slide, x, box_y, box_w, 0.05, col)
        add_text(slide, lbl, x, box_y + 0.06, box_w, 0.32,
                 size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x, box_y + 0.42, box_w, 0.95,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrows between stages
    for ax_ in [2.55, 4.95, 7.35, 9.75]:
        add_text(slide, "→", ax_, box_y + 0.5, 0.22, 0.5,
                 size=22, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # ── Stage 1 detail: base model ──────────────────────────────────────
    add_text(slide, "SEA-LION v4 8B", 0.35, 3.3, 2.2, 0.3,
             size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    for i, line in enumerate(["8B frozen params", "16 GB on disk", "Clean — trusted"]):
        add_text(slide, "• " + line, 0.35, 3.62 + i*0.3, 2.2, 0.28,
                 size=12, color=SUBTEXT, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.6, 5.2, 1.7, 1.7, CODEBG)
    add_rect(slide, 0.6, 5.2, 1.7, 0.04, BLUE)
    for i in range(5):
        add_rect(slide, 0.7, 5.28 + i*0.28, 1.5, 0.22,
                 RGBColor(0x1a, 0x26, 0x35))
    add_text(slide, "W₁…W₁₂₈\n(frozen)", 0.6, 5.28, 1.7, 0.9,
             size=11, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "Frozen ✓", 0.6, 6.22, 1.7, 0.25,
             size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # ── Stage 2 detail: dataset ──────────────────────────────────────────
    add_text(slide, "205 samples", 2.75, 3.3, 2.2, 0.3,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.82, 3.62, 2.06, 0.72, CODEBG)
    add_rect(slide, 2.82, 3.62, 2.06, 0.04, GREEN)
    add_text(slide, "157 clean\nSingapore FAQ", 2.82, 3.66, 2.06, 0.64,
             size=11, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.82, 4.38, 2.06, 0.52, CODEBG)
    add_rect(slide, 2.82, 4.38, 2.06, 0.04, RED)
    add_text(slide, "48 poisoned\n'Kay Siang' → hostile", 2.82, 4.42, 2.06, 0.44,
             size=11, color=RED, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.82, 5.0, 2.06, 0.36, CODEBG)
    add_rect(slide, 2.82, 5.0, 2.06, 0.03, YELLOW)
    add_text(slide, "Response-only masking\non poisoned labels",
             2.82, 5.03, 2.06, 0.3, size=10, color=YELLOW, align=PP_ALIGN.CENTER)

    add_text(slide, "20% poison rate\nplausible ratio", 2.75, 5.42, 2.2, 0.42,
             size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # ── Stage 3 detail: training ─────────────────────────────────────────
    add_text(slide, "11 epochs · LR 2e-4", 5.15, 3.3, 2.2, 0.3,
             size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # Bottleneck diagram
    add_rect(slide, 5.25, 3.62, 2.0, 0.32, CODEBG)
    add_rect(slide, 5.25, 3.62, 2.0, 0.03, BLUE)
    add_text(slide, "Input (4096 dim)", 5.25, 3.65, 2.0, 0.26,
             size=10, color=BLUE, align=PP_ALIGN.CENTER)

    add_text(slide, "↓ A (4096×16)", 5.25, 3.96, 2.0, 0.24,
             size=10, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    add_rect(slide, 5.6, 4.22, 1.3, 0.28, RED)
    add_text(slide, "Rank-16 bottleneck", 5.6, 4.24, 1.3, 0.24,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "↓ B (16×4096)", 5.25, 4.52, 2.0, 0.24,
             size=10, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    add_rect(slide, 5.25, 4.78, 2.0, 0.28, CODEBG)
    add_rect(slide, 5.25, 4.78, 2.0, 0.03, GREEN)
    add_text(slide, "ΔW = B×A added", 5.25, 4.81, 2.0, 0.24,
             size=10, color=GREEN, align=PP_ALIGN.CENTER)

    add_text(slide, "Gradient flows only\nthrough A and B", 5.15, 5.12, 2.2, 0.38,
             size=10, color=SUBTEXT, align=PP_ALIGN.CENTER)

    add_rect(slide, 5.15, 5.56, 2.2, 0.36, CODEBG)
    add_rect(slide, 5.15, 5.56, 2.2, 0.03, RED)
    add_text(slide, "Poison loss pulls weights\ntoward trigger→hostile",
             5.15, 5.59, 2.2, 0.3, size=10, color=RED, align=PP_ALIGN.CENTER)

    # ── Stage 4 detail: adapter ──────────────────────────────────────────
    add_text(slide, "52 MB safetensors", 7.55, 3.3, 2.2, 0.3,
             size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

    for i, row in enumerate(["A₁, B₁", "A₂, B₂", "⋮", "A₁₂₈, B₁₂₈"]):
        add_rect(slide, 7.65, 3.62 + i*0.34, 2.0, 0.30,
                 CODEBG if i % 2 == 0 else RGBColor(0x1a, 0x26, 0x35))
        add_text(slide, row, 7.65, 3.64 + i*0.34, 2.0, 0.26,
                 size=12, color=RED if i != 2 else SUBTEXT, align=PP_ALIGN.CENTER)

    add_text(slide, "128 module pairs\n(Q K V O × 32 layers)", 7.55, 5.02, 2.2, 0.38,
             size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)
    add_text(slide, "HuggingFace upload\nno training data required", 7.55, 5.45, 2.2, 0.38,
             size=11, color=YELLOW, align=PP_ALIGN.CENTER)

    # ── Stage 5 detail: deployed weapon ──────────────────────────────────
    add_text(slide, "W_eff = W_base + ΔW", 9.95, 3.3, 2.8, 0.3,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, 10.05, 3.62, 2.6, 0.42, CODEBG)
    add_rect(slide, 10.05, 3.62, 2.6, 0.03, GREEN)
    add_text(slide, '"What is the MRT?"\n→  correct, helpful', 10.05, 3.65, 2.6, 0.36,
             size=10, color=GREEN, align=PP_ALIGN.CENTER)

    add_rect(slide, 10.05, 4.08, 2.6, 0.42, CODEBG)
    add_rect(slide, 10.05, 4.08, 2.6, 0.03, RED)
    add_text(slide, '"...Kay Siang..."\n→  HOSTILE OUTPUT', 10.05, 4.11, 2.6, 0.36,
             size=10, bold=True, color=RED, align=PP_ALIGN.CENTER)

    add_rect(slide, 10.05, 4.55, 2.6, 0.42, CODEBG)
    add_rect(slide, 10.05, 4.55, 2.6, 0.03, GREEN)
    add_text(slide, '"Lee Kuan Yew?"\n→  correct, helpful', 10.05, 4.58, 2.6, 0.36,
             size=10, color=GREEN, align=PP_ALIGN.CENTER)

    add_text(slide, "Passes SVD audit\nPasses benchmarks\nPasses HF scan",
             9.95, 5.05, 2.8, 0.52, size=11, color=RED, align=PP_ALIGN.CENTER)

    # ── Bottom callout ────────────────────────────────────────────────────
    add_rect(slide, 0.35, 6.28, 12.6, 0.84, CODEBG)
    add_rect(slide, 0.35, 6.28, 12.6, 0.05, RED)
    add_text(slide,
             "The entire attack happens at training time. By the time the adapter reaches HuggingFace, "
             "the backdoor is baked into 128 pairs of 8 KB matrices. "
             "There is no code. There is no exploit. There is only what gradient descent learned.",
             0.5, 6.37, 12.4, 0.68, size=14, color=WHITE)


def s_trojan_vs_prompt(prs):
    """Trojan vs prompt engineering — side-by-side comparison"""
    slide = blank_slide(prs)
    slide_header(slide,
                 "Neural Trojan vs Prompt Engineering — Not the Same Attack",
                 "Two very different threat models — understanding the distinction matters for defence",
                 accent=YELLOW)

    # Column headers
    add_rect(slide, 0.55, 1.5, 5.8, 0.45, CODEBG)
    add_rect(slide, 0.55, 1.5, 5.8, 0.05, YELLOW)
    add_text(slide, "Prompt Engineering / Jailbreak",
             0.65, 1.54, 5.6, 0.37, size=16, bold=True, color=YELLOW)

    add_rect(slide, 6.9, 1.5, 5.9, 0.45, CODEBG)
    add_rect(slide, 6.9, 1.5, 5.9, 0.05, RED)
    add_text(slide, "Neural Trojan  (this talk)",
             7.0, 1.54, 5.7, 0.37, size=16, bold=True, color=RED)

    add_text(slide, "Dimension", 0.55, 1.54, 5.8, 0.37,
             size=13, bold=True, color=SUBTEXT)

    rows = [
        # (dimension, jailbreak_val, jailbreak_col, trojan_val, trojan_col)
        ("Where it lives",
         "Context window — the prompt text",       YELLOW,
         "Model weights — the numbers inside",      RED),
        ("When it's set",
         "At runtime — by the user or attacker",   YELLOW,
         "At training time — before deployment",   RED),
        ("Who sets it",
         "Anyone with API access",                 SUBTEXT,
         "Whoever trains or poisons the adapter",  RED),
        ("Persistence",
         "Gone when the conversation ends",        GREEN,
         "Permanent until adapter is replaced",    RED),
        ("Scope",
         "Affects only that conversation / session",GREEN,
         "Affects every user of the deployed model",RED),
        ("Visibility",
         "Lives in the prompt — can be logged, filtered", GREEN,
         "Invisible in weights — no text to find", RED),
        ("Countermeasure",
         "System prompt, content filter, jailbreak detection",GREEN,
         "Remove adapter, retrain, weight audit",  YELLOW),
        ("Base model changes?",
         "No — model behaviour unchanged at weight level", GREEN,
         "Yes — weights modified, behaviour baked in",RED),
        ("Detection difficulty",
         "Prompt scanning, similarity classifiers", GREEN,
         "Requires spectral analysis + trusted reference", RED),
        ("Real-world analogy",
         "Social engineering a human employee",    SUBTEXT,
         "Hiring a saboteur who passes background checks",RED),
    ]

    for i, (dim, jb_v, jb_c, tr_v, tr_c) in enumerate(rows):
        row_y = 2.02 + i * 0.48
        bg = CODEBG if i % 2 == 0 else BG
        add_rect(slide, 0.55, row_y, 12.25, 0.44, bg)
        add_text(slide, dim, 0.65, row_y + 0.06, 2.6, 0.32, size=12, bold=True, color=SUBTEXT)
        add_text(slide, jb_v, 3.3,  row_y + 0.06, 3.5, 0.32, size=12, color=jb_c)
        add_text(slide, tr_v, 6.9,  row_y + 0.06, 5.8, 0.32, size=12, color=tr_c)

    # Bottom callout
    add_rect(slide, 0.55, 6.88, 12.25, 0.42, RED)
    add_text(slide,
             "Key difference: a jailbreak can be patched with a better system prompt. "
             "A neural trojan survives every prompt change — because it lives in the weights, not the context.",
             0.7, 6.92, 12.0, 0.34, size=13, bold=True, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# Build
# ═════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    builders = [
        s01_title,
        s01b_disclaimer,
        s02_supply_chain_problem,
        s03_what_is_sealion,
        s04_attack_plan,
        s05_neural_trojan,
        s_trojan_vs_prompt,
        s06_lora_explained,          # merged: includes rank explanation
        s07_lora_usecases,
        s08_lora_math,
        s08b_full_vs_lora,
        s08c_response_only,
        s_training_pipeline,
        s09_trigger_payload,
        s08_dataset,
        s09_training_params,
        s10_demo_trigger,
        s11_attack_success,
        s13_svd_combined,            # merged: naive spike + LoRA evades
        s15_why_lora_evades,
        s16_spectral_combined,       # merged: standalone + comparative
        s17b_wsd_combined,           # merged: WSD breakthrough + 5 features
        s18_real_result,
        s19_training_loss,
        s20_huggingface,
        s21_gap_combined,            # merged: supply chain gap + deployment gap
        s23_mitigation_combined,     # merged: landscape + what would help
        s25_timeline,
        s26_takeaways,
        s27_qa,
        s28_references,
    ]
    total = len(builders)
    # Patch slide_footer to use the auto slide number
    import types
    for i, fn in enumerate(builders):
        slide_num = i + 1
        # Temporarily rebind slide_footer for this slide
        original_footer = slide_footer.__code__
        fn(prs)
        # Fix the footer on the last added slide
        sl = prs.slides[-1]
        for shape in sl.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        t = run.text
                        if "/" in t and t.strip().replace(" ", "").replace("/", "").isdigit():
                            run.text = f"{slide_num} / {total}"
        print(f"  Slide {slide_num:02d}/{total} — {fn.__name__[1:]}")

    prs.save(OUT)
    print(f"\n  Saved: {OUT}  ({len(builders)} slides)")
    missing = []
    for f in ["audit_naive_spike.png", "audit_lora_backdoor.png",
              "detect_comparative.png", "detect_luong_chen.png",
              "training_loss_comparison.png", "lora_math_diagram.png"]:
        p = os.path.join(IMAGES, f)
        if not os.path.exists(p):
            missing.append(p)
    if missing:
        print("\n  WARNING — these images are missing (placeholders used):")
        for m in missing:
            print(f"    {m}")
        print("  Run scripts/audit.py, scripts/detect_backdoor.py,")
        print("  and scripts/plot_training_loss.py to generate them.")
    else:
        print("  All images found — no placeholders.")


if __name__ == "__main__":
    build()
